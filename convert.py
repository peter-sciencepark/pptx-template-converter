#!/usr/bin/env python3
"""
pptx-template-converter

Converts PowerPoint presentations from an old layout/template
to a new one, mapping slide types and transferring content.

Usage:
    python3 convert.py input.pptx --template template.potx -o output.pptx
"""

import argparse
import io
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN


# ---------------------------------------------------------------------------
# Template handling
# ---------------------------------------------------------------------------

def open_potx_as_presentation(potx_path: str) -> Presentation:
    """Open a .potx template file as a Presentation by patching the content type."""
    buf = io.BytesIO()
    with zipfile.ZipFile(potx_path, "r") as z_in, zipfile.ZipFile(buf, "w") as z_out:
        for item in z_in.infolist():
            data = z_in.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = data.replace(
                    b"application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
                    b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
                )
            z_out.writestr(item, data)
    buf.seek(0)
    return Presentation(buf)


def get_layout_map(template_prs: Presentation) -> dict:
    """Build a name -> layout mapping from the template."""
    layouts = {}
    for master in template_prs.slide_masters:
        for layout in master.slide_layouts:
            layouts[layout.name] = layout
    return layouts


# ---------------------------------------------------------------------------
# Text extraction helpers
# ---------------------------------------------------------------------------

def _safe_color(font):
    """Safely extract RGB color string from a font, or return None."""
    try:
        if font.color and font.color.type is not None:
            return str(font.color.rgb)
    except AttributeError:
        pass
    return None


def extract_all_text(shape) -> list[dict]:
    """Extract text runs from a shape, preserving paragraph structure."""
    paragraphs = []
    if not shape.has_text_frame:
        return paragraphs
    for para in shape.text_frame.paragraphs:
        runs = []
        for run in para.runs:
            runs.append({
                "text": run.text,
                "bold": run.font.bold,
                "italic": run.font.italic,
                "size": run.font.size,
                "color": _safe_color(run.font),
            })
        if not runs and para.text:
            runs.append({"text": para.text, "bold": None, "italic": None, "size": None, "color": None})
        paragraphs.append({
            "runs": runs,
            "alignment": para.alignment,
            "level": para.level,
        })
    return paragraphs


# ---------------------------------------------------------------------------
# Slide content extraction
# ---------------------------------------------------------------------------

def find_title_text(slide) -> str:
    """Find the most likely title text from a slide."""
    # Check placeholders first
    for shape in slide.placeholders:
        if shape.placeholder_format.type in (1, 3):  # TITLE, CENTER_TITLE
            return shape.text.strip()

    # Check named shapes — look for the largest/topmost text shape as title candidate
    for shape in slide.shapes:
        name = shape.name.lower()
        if "rubrik" in name or "title" in name or "titel" in name:
            if hasattr(shape, "text") and shape.text.strip():
                return shape.text.strip()

    # Heuristic: if first text shape is short and near top, treat as title
    text_shapes = [s for s in slide.shapes
                   if hasattr(s, "text") and s.text.strip() and s.shape_type != 13]
    if text_shapes:
        # Sort by vertical position
        text_shapes.sort(key=lambda s: s.top)
        candidate = text_shapes[0]
        words = len(candidate.text.split())
        # Short text near the top of the slide
        if words <= 10 and candidate.top < Inches(2):
            return candidate.text.strip()

    return ""


def find_body_text(slide) -> list[dict]:
    """Find the main body/content text from a slide (excluding title and sources)."""
    title_text = find_title_text(slide)
    source_text = find_source_text(slide)
    all_paragraphs = []

    for shape in slide.shapes:
        if not hasattr(shape, "text") or not shape.text.strip():
            continue
        if shape.shape_type == 13:  # Skip pictures
            continue
        # Skip if this is the title
        if shape.text.strip() == title_text and title_text:
            continue
        # Skip source/attribution text
        if source_text and shape.text.strip() == source_text:
            continue

        paragraphs = extract_all_text(shape)
        all_paragraphs.extend(paragraphs)

    return all_paragraphs


def find_source_text(slide) -> str:
    """Find source/attribution text (e.g. 'Källa: ...')."""
    for shape in slide.shapes:
        if not hasattr(shape, "text") or not shape.text.strip():
            continue
        text_lower = shape.text.lower().strip()
        if any(kw in text_lower for kw in ["källa", "source", "ref:", "referens"]):
            return shape.text.strip()
    return ""


def find_images(slide) -> list:
    """Find all picture shapes in a slide."""
    return [s for s in slide.shapes if s.shape_type == 13]


# ---------------------------------------------------------------------------
# Slide classification
# ---------------------------------------------------------------------------

def classify_slide(slide, index: int, total: int) -> str:
    """
    Classify a slide into a category based on its content and position.

    Returns one of:
        'title'       - Title/cover slide
        'chapter'     - Chapter heading (short text, no bullets)
        'quote'       - Quote slide
        'content'     - Content with text/bullets
        'content_img' - Content with image
        'closing'     - Last slide / closing (only if minimal content)
        'blank'       - Empty or minimal content
    """
    layout_name = slide.slide_layout.name.lower()
    shapes = list(slide.shapes)

    has_picture = any(s.shape_type == 13 for s in shapes)
    text_shapes = [s for s in shapes if hasattr(s, "text") and s.text.strip()]
    all_text = " ".join(s.text for s in text_shapes).strip()
    word_count = len(all_text.split()) if all_text else 0

    # First slide is typically a title
    if index == 0:
        return "title"

    # Last slide — only classify as closing if very little content
    if index == total - 1 and word_count < 10 and not has_picture:
        return "closing"

    # Quote detection: starts with quotation marks
    if all_text.startswith('"') or all_text.startswith('\u201c') or all_text.startswith('\u201d'):
        return "quote"
    if "citat" in layout_name:
        return "quote"

    # Blank
    if word_count == 0 and not has_picture:
        return "blank"

    # Content with image
    if has_picture:
        return "content_img"

    # Title-like slides with very little text and no content shapes
    if "rubrik" in layout_name and word_count < 12:
        return "chapter"

    # Short text without bullets = chapter
    if word_count < 10 and len(text_shapes) <= 2:
        return "chapter"

    return "content"


# ---------------------------------------------------------------------------
# Layout mapping
# ---------------------------------------------------------------------------

def map_to_new_layout(category: str, layouts: dict) -> str:
    """Map a slide category to the best matching new template layout name.

    Strategy: prefer layouts WITH placeholders so content goes into the
    template's styled text frames rather than plain textboxes.
    """
    mapping = {
        "title": ["1 - Rubrikbild logo", "1 - Rubrikbild blank"],
        "chapter": [
            "1 - Kapitelrubrik med underrubrik",
            "1 - Kapitelrubrik",
        ],
        "quote": ["11 - Midicitat blå", "11 - Maxicitat blå"],
        "content": [
            # "Bild höger" has TITLE + OBJECT placeholders — use it even
            # without an image so text lands in styled placeholders.
            "4 - Bild höger",
            "5 - Bild höger",
            "7 - Bild höger",
            "4 - Innehåll blank",
        ],
        "content_img": [
            "4 - Bild höger",
            "5 - Bild höger",
            "7 - Bild höger",
        ],
        "closing": [
            "13 - Bakgrund hav",
            "14 - Slogan hav",
            "13 - Bakgrund skog",
        ],
        "blank": [
            "4 - Innehåll blank",
            "5 - Innehåll blank",
        ],
    }

    candidates = mapping.get(category, ["4 - Bild höger"])
    for name in candidates:
        if name in layouts:
            return name

    # Fallback
    for name in layouts:
        if "Bild höger" in name:
            return name
    return list(layouts.keys())[0]


# ---------------------------------------------------------------------------
# Writing content to slides
# ---------------------------------------------------------------------------

def apply_text_to_textframe(text_frame, paragraphs: list[dict]):
    """Write extracted paragraph data into a text frame."""
    text_frame.clear()

    first_para = True
    for para_data in paragraphs:
        if not para_data["runs"]:
            continue

        if first_para:
            p = text_frame.paragraphs[0]
            first_para = False
        else:
            p = text_frame.add_paragraph()

        if para_data["alignment"] is not None:
            p.alignment = para_data["alignment"]
        p.level = para_data.get("level", 0) or 0

        for j, run_data in enumerate(para_data["runs"]):
            if j == 0 and p.runs:
                r = p.runs[0]
            else:
                r = p.add_run()
            r.text = run_data["text"]
            if run_data["bold"] is not None:
                r.font.bold = run_data["bold"]
            if run_data["italic"] is not None:
                r.font.italic = run_data["italic"]


def add_source_text(slide, source_text: str):
    """Add a small source/attribution text at the bottom of the slide."""
    if not source_text:
        return
    txBox = slide.shapes.add_textbox(
        Inches(7.5), Inches(6.6), Inches(5.0), Inches(0.4)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = source_text
    run.font.size = Pt(9)
    run.font.italic = True
    run.font.color.rgb = _parse_rgb("999999")


def _parse_rgb(hex_str):
    """Convert hex string to RGBColor."""
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(hex_str)


# ---------------------------------------------------------------------------
# Slide conversion
# ---------------------------------------------------------------------------

def convert_slide(old_slide, new_prs, layouts, category, slide_index):
    """Convert a single slide from old to new template."""
    layout_name = map_to_new_layout(category, layouts)
    layout = layouts[layout_name]
    new_slide = new_prs.slides.add_slide(layout)

    title_text = find_title_text(old_slide)
    body_paragraphs = find_body_text(old_slide)
    source_text = find_source_text(old_slide)
    images = find_images(old_slide)

    # Collect available placeholders from the new slide
    title_ph = None
    body_ph = None
    subtitle_ph = None
    picture_ph = None

    for ph in new_slide.placeholders:
        pt = ph.placeholder_format.type
        idx = ph.placeholder_format.idx
        if pt in (1, 3):  # TITLE, CENTER_TITLE
            title_ph = ph
        elif pt == 18:  # PICTURE
            if picture_ph is None:
                picture_ph = ph
        elif pt == 4:  # SUBTITLE
            subtitle_ph = ph
        elif pt in (2, 7):  # BODY, OBJECT
            if body_ph is None:
                body_ph = ph

    # --- Chapter slides: special handling ---
    # "Kapitelrubrik" has only BODY (idx=10) and optionally SUBTITLE (idx=1).
    # Put the title into the BODY placeholder (which is styled as the main heading).
    if category == "chapter":
        if body_ph and title_text:
            body_ph.text = title_text
        elif title_text:
            _add_title_textbox(new_slide, title_text)
        if subtitle_ph and body_paragraphs:
            apply_text_to_textframe(subtitle_ph.text_frame, body_paragraphs)
        add_source_text(new_slide, source_text)
        return new_slide, layout_name

    # --- Title/cover slides: minimal content ---
    if category == "title":
        # Title slides in the new template have no placeholders — the logo
        # and background come from the master.  Nothing to fill.
        add_source_text(new_slide, source_text)
        return new_slide, layout_name

    # --- Closing slides ---
    if category == "closing":
        add_source_text(new_slide, source_text)
        return new_slide, layout_name

    # --- Quote slides ---
    if category == "quote":
        if body_ph:
            # Put all text (including the quote) into the body placeholder
            all_paras = []
            if title_text:
                all_paras.append({
                    "runs": [{"text": title_text, "bold": True, "italic": None, "size": None, "color": None}],
                    "alignment": None, "level": 0,
                })
            all_paras.extend(body_paragraphs)
            if all_paras:
                apply_text_to_textframe(body_ph.text_frame, all_paras)
            elif title_text:
                body_ph.text = title_text
        add_source_text(new_slide, source_text)
        return new_slide, layout_name

    # --- Content slides (with or without image) ---

    # Apply title
    if title_text and title_ph:
        title_ph.text = title_text
    elif title_text and not title_ph:
        _add_title_textbox(new_slide, title_text)

    # Apply body content
    if body_paragraphs and body_ph:
        apply_text_to_textframe(body_ph.text_frame, body_paragraphs)
    elif body_paragraphs and not body_ph:
        top = Inches(2.1) if title_text else Inches(1.5)
        txBox = new_slide.shapes.add_textbox(
            Inches(0.92), top, Inches(11.5), Inches(4.5)
        )
        apply_text_to_textframe(txBox.text_frame, body_paragraphs)

    # Handle images
    if images and picture_ph:
        _insert_image_to_placeholder(images[0], picture_ph, new_slide)
    elif images:
        _add_image_as_shape(new_slide, images)

    # Source attribution
    add_source_text(new_slide, source_text)

    # If this is a content slide using "Bild höger" but with no image,
    # remove the empty picture placeholder to avoid a blank box
    if not images and picture_ph and category == "content":
        _remove_placeholder(new_slide, picture_ph)

    return new_slide, layout_name


def _add_title_textbox(slide, title_text: str):
    """Add a styled title as a textbox when no title placeholder exists."""
    txBox = slide.shapes.add_textbox(
        Inches(0.92), Inches(1.02), Inches(11.5), Inches(0.83)
    )
    p = txBox.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = _parse_rgb("0C4C4B")


def _insert_image_to_placeholder(image_shape, placeholder, slide=None):
    """Insert an image into a picture placeholder."""
    try:
        image_stream = io.BytesIO(image_shape.image.blob)
        placeholder.insert_picture(image_stream)
    except Exception as e:
        print(f"    Warning: Could not insert image into placeholder: {e}")
        # Fallback: add as free shape at the placeholder's position
        if slide:
            try:
                image_stream = io.BytesIO(image_shape.image.blob)
                slide.shapes.add_picture(
                    image_stream,
                    placeholder.left, placeholder.top,
                    placeholder.width, placeholder.height,
                )
                print(f"    -> Added image as free shape instead")
            except Exception as e2:
                print(f"    -> Fallback also failed: {e2}")


def _add_image_as_shape(slide, images):
    """Add images as free-floating shapes on the slide."""
    for img in images:
        try:
            image_stream = io.BytesIO(img.image.blob)
            slide.shapes.add_picture(
                image_stream,
                img.left, img.top,
                img.width, img.height,
            )
        except Exception as e:
            print(f"    Warning: Could not add image: {e}")


def _remove_placeholder(slide, placeholder):
    """Remove an empty placeholder element from a slide."""
    try:
        sp = placeholder._element
        sp.getparent().remove(sp)
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Main conversion
# ---------------------------------------------------------------------------

def convert_presentation(input_path: str, template_path: str, output_path: str):
    """Main conversion function."""
    print(f"Opening input:    {input_path}")
    old_prs = Presentation(input_path)

    print(f"Opening template: {template_path}")
    if template_path.endswith(".potx"):
        new_prs = open_potx_as_presentation(template_path)
    else:
        new_prs = Presentation(template_path)

    # Remove example slides from template
    ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    while len(new_prs.slides) > 0:
        sldId = new_prs.slides._sldIdLst[0]
        rId = sldId.get(f"{{{ns_r}}}id") or sldId.get("r:id")
        if rId:
            new_prs.part.drop_rel(rId)
        new_prs.slides._sldIdLst.remove(sldId)

    layouts = get_layout_map(new_prs)
    print(f"\nAvailable layouts in new template:")
    for name in sorted(layouts.keys()):
        phs = list(layouts[name].placeholders)
        ph_info = ", ".join(f"{p.placeholder_format.type}" for p in phs) if phs else "no placeholders"
        print(f"  - {name}  [{ph_info}]")

    total_slides = len(old_prs.slides)
    print(f"\nConverting {total_slides} slides...\n")

    for i, slide in enumerate(old_prs.slides):
        category = classify_slide(slide, i, total_slides)
        old_layout = slide.slide_layout.name
        new_slide, new_layout = convert_slide(slide, new_prs, layouts, category, i)

        # Log details
        title = find_title_text(slide)
        imgs = len(find_images(slide))
        src = find_source_text(slide)
        details = []
        if title:
            details.append(f'title="{title[:30]}"')
        if imgs:
            details.append(f"images={imgs}")
        if src:
            details.append(f'source="{src[:25]}"')
        detail_str = f"  ({', '.join(details)})" if details else ""
        print(f"  Slide {i+1}/{total_slides}: [{category:12s}] '{old_layout}' -> '{new_layout}'{detail_str}")

    new_prs.save(output_path)
    print(f"\nSaved: {output_path}")


def main():
    parser = argparse.ArgumentParser(
        description="Convert PowerPoint presentations to a new template layout."
    )
    parser.add_argument("input", help="Input .pptx file to convert")
    parser.add_argument(
        "--template", "-t", required=True,
        help="New template file (.potx or .pptx)"
    )
    parser.add_argument(
        "--output", "-o", default=None,
        help="Output .pptx file (default: input_ny_mall.pptx)"
    )
    args = parser.parse_args()

    if args.output is None:
        stem = Path(args.input).stem
        args.output = str(Path(args.input).parent / f"{stem}_ny_mall.pptx")

    convert_presentation(args.input, args.template, args.output)


if __name__ == "__main__":
    main()
